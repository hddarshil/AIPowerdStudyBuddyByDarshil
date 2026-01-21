import streamlit as st
import google.generativeai as genai
import toml
import PyPDF2
from pptx import Presentation
import speech_recognition as sr
import tempfile
from datetime import date
from io import BytesIO
import time
from fpdf import FPDF
from PIL import Image, ImageDraw, ImageFont
import networkx as nx
import matplotlib.pyplot as plt

# ==================================================
# PAGE CONFIG
# ==================================================
st.set_page_config(
    page_title="AI Study Buddy Pro 🤖",
    layout="wide",
    page_icon="📘"
)

# ==================================================
# SESSION STATE
# ==================================================
if "theme" not in st.session_state:
    st.session_state.theme="light"
if "history" not in st.session_state:
    st.session_state.history=[]
if "activity" not in st.session_state:
    st.session_state.activity=0

# ==================================================
# THEME TOGGLE
# ==================================================
def toggle_theme():
    st.session_state.theme="dark" if st.session_state.theme=="light" else "light"
st.sidebar.button(f"🌙 Toggle Dark/Light Theme", on_click=toggle_theme)

# ==================================================
# STYLE + ANIMATION
# ==================================================
def apply_style():
    if st.session_state.theme=="light":
        st.markdown("""
        <style>
        [data-testid="stAppViewContainer"]{
            background: linear-gradient(135deg,#e3f2fd,#fce4ec,#e8f5e9);
            background-size:400% 400%;
            animation: gradient 20s ease infinite;
            transition: background 0.5s ease;
        }
        @keyframes gradient{
            0%{background-position:0% 50%;}
            50%{background-position:100% 50%;}
            100%{background-position:0% 50%;}
        }
        .glass{background:rgba(255,255,255,0.88);border-radius:22px;padding:28px;
        box-shadow:0 12px 40px rgba(0,0,0,0.12);backdrop-filter:blur(18px);margin-bottom:25px;}
        .chat-user{background:#42a5f5;color:white;padding:14px;border-radius:18px;margin:8px 0;text-align:right;animation:fadein 0.6s;}
        .chat-ai{background:#f5f5f5;padding:14px;border-radius:18px;margin:8px 0;border-left:5px solid #7e57c2;animation:fadein 0.6s;}
        @keyframes fadein{0%{opacity:0;}100%{opacity:1;}}
        .stButton>button{background:linear-gradient(135deg,#42a5f5,#7e57c2);color:white;border-radius:14px;padding:12px;font-weight:bold;border:none;transition:0.3s;}
        .stButton>button:hover{transform:scale(1.05);background:linear-gradient(135deg,#66bb6a,#26c6da);}
        </style>
        """, unsafe_allow_html=True)
    else:
        st.markdown("""
        <style>
        [data-testid="stAppViewContainer"]{background:#121212;color:white;transition:0.5s;}
        .glass{background:rgba(0,0,0,0.7);border-radius:22px;padding:28px;
        box-shadow:0 12px 40px rgba(0,0,0,0.5);backdrop-filter:blur(18px);margin-bottom:25px;color:white;}
        .chat-user{background:#1e88e5;color:white;padding:14px;border-radius:18px;margin:8px 0;text-align:right;animation:fadein 0.6s;}
        .chat-ai{background:#333;padding:14px;border-radius:18px;margin:8px 0;border-left:5px solid #7e57c2;color:white;animation:fadein 0.6s;}
        .stButton>button{background:linear-gradient(135deg,#42a5f5,#7e57c2);color:white;border-radius:14px;padding:12px;font-weight:bold;border:none;transition:0.3s;}
        .stButton>button:hover{transform:scale(1.05);background:linear-gradient(135deg,#66bb6a,#26c6da);}
        @keyframes fadein{0%{opacity:0;}100%{opacity:1;}}
        </style>
        """, unsafe_allow_html=True)
apply_style()

# ==================================================
# LOAD GEMINI API
# ==================================================
try:
    config = toml.load("config.toml")
    API_KEY = config["gemini"]["api_key"]
except:
    API_KEY = None

model = None
if API_KEY:
    genai.configure(api_key=API_KEY)
    model = genai.GenerativeModel("gemini-2.5-flash")

# ==================================================
# GEMINI FUNCTION
# ==================================================
def ask_gemini(prompt):
    try:
        response = model.generate_content(f"Answer concisely for students:\n{prompt}")
        return response.text
    except:
        return "⚠️ AI response error."

# ==================================================
# UNIVERSAL INPUT
# ==================================================
def get_input_text():
    text = st.text_area("✍️ Enter text (optional)")
    file = st.file_uploader("📂 Upload file (PDF / TXT / PPT)", type=["pdf","txt","pptx"])
    if file:
        if file.type=="application/pdf":
            reader = PyPDF2.PdfReader(file)
            return "".join(p.extract_text() for p in reader.pages[:5])
        elif file.type=="text/plain":
            return file.read().decode()
        elif file.type=="application/vnd.openxmlformats-officedocument.presentationml.presentation":
            ppt = Presentation(file)
            data = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape,"text"):
                        data+=shape.text
            return data
    return text

# ==================================================
# SIDEBAR MENU
# ==================================================
st.sidebar.title("🎓 AI Study Buddy Pro")
menu = st.sidebar.radio("Menu",[
    "🏠 Home","🧠 Explain Topic","📝 Summarize Notes","📄 PDF Summarizer","🎯 Quiz Generator",
    "💡 Flashcards","🎤 Voice Question","📊 Notes → PPT","🧠 Smart Study Planner ⭐",
    "📈 Knowledge Analyzer ⭐","🔁 Smart Revision Mode ⭐","⏱ Focus Booster ⭐",
    "📊 Learning Tracker ⭐","📂 History","📄 Export PDF","🖼 Concept Map","🎮 Memory Challenge","ℹ️ About"
])

# ==================================================
# CHAT SECTION WITH TYPING
# ==================================================
def chat_section(title,prompt_template):
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    st.subheader(title)
    content=get_input_text()
    if st.button("🚀 Generate"):
        st.markdown(f"<div class='chat-user'>{content}</div>",unsafe_allow_html=True)
        placeholder=st.empty()
        text=""
        ans=ask_gemini(prompt_template+content)
        for c in ans:
            text+=c
            placeholder.markdown(f"<div class='chat-ai'>{text}</div>",unsafe_allow_html=True)
            time.sleep(0.01)
        st.session_state.history.append((title,content,ans))
        st.session_state.activity+=1
    st.markdown("</div>",unsafe_allow_html=True)

# ==================================================
# FEATURES IMPLEMENTATION
# ==================================================
# if menu=="🏠 Home":
#     st.markdown("<div class='glass'>",unsafe_allow_html=True)
#     st.title("🤖 Neo Learn Buddy Ai Pro")
#     st.write("### Modern • Animated • Interactive Learning Assistant 🚀")
#     st.markdown("</div>",unsafe_allow_html=True)
if menu=="🏠 Home":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    st.title("🤖 NeoLearn Buddy AI")
    st.write("### Modern • Animated • Interactive Learning Assistant 🚀")
    st.markdown("""
### Welcome to **NeoLearn Buddy AI** 🤖📘

Your ultimate **AI-powered study companion** designed for college and diploma students.  

**What it can do for you:**  
✔ Explain complex topics in simple language  
✔ Summarize your notes, PDFs, and presentations  
✔ Generate quizzes and flashcards for effective learning  
✔ Convert notes into interactive PPTs  
✔ Analyze your knowledge level and suggest study strategies  
✔ Create concept maps and diagrams for visual learning  
✔ Memory challenge games to boost recall  
✔ Focus booster and study tracker to maintain consistency  

**Why NeoLearn Buddy AI?**  
- Modern, animated UI/UX for an engaging experience  
- All-in-one study assistant — no need for multiple apps  
- Personalized and interactive — adapts to your learning style  
- Save, export, and share your study material with ease  

Start learning smarter and faster with **NeoLearn Buddy AI**! 🚀
""")
    st.markdown("</div>",unsafe_allow_html=True)


elif menu=="🧠 Explain Topic":
    chat_section("Explain Topic","Explain in simple words:\n")

elif menu=="📝 Summarize Notes":
    chat_section("Summarize Notes","Summarize into bullet points:\n")

elif menu=="📄 PDF Summarizer":
    chat_section("PDF Summary","Summarize student-friendly:\n")

elif menu=="🎯 Quiz Generator":
    chat_section("Quiz Generator","Create 5 MCQs with answers:\n")

elif menu=="💡 Flashcards":
    chat_section("Flashcards","Create flashcards term:definition:\n")

elif menu=="🔁 Smart Revision Mode ⭐":
    chat_section("Revision Mode","Exam-focused revision notes:\n")

elif menu=="📈 Knowledge Analyzer ⭐":
    chat_section("Knowledge Analyzer","Analyze knowledge level and tips:\n")

elif menu=="⏱ Focus Booster ⭐":
    chat_section("Focus Booster","Give motivation for focused study:\n")

elif menu=="📊 Notes → PPT":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    notes=get_input_text()
    if st.button("Generate PPT"):
        ppt=Presentation()
        slide=ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text="AI Notes"
        slide.placeholders[1].text=ask_gemini(notes)
        ppt.save("notes.pptx")
        st.download_button("⬇️ Download PPT","notes.pptx")
    st.markdown("</div>",unsafe_allow_html=True)

elif menu=="🧠 Smart Study Planner ⭐":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    subject=st.text_input("Subject")
    hours=st.slider("Study hours per day",1,10)
    exam=st.date_input("Exam date",date.today())
    if st.button("Generate Plan"):
        ans=ask_gemini(f"Create a daily study plan for {subject} with {hours} hours until {exam}.")
        st.write(ans)
    st.markdown("</div>",unsafe_allow_html=True)

elif menu=="🎤 Voice Question":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    audio=st.file_uploader("Upload voice (wav)",type=["wav"])
    if audio:
        r=sr.Recognizer()
        with tempfile.NamedTemporaryFile(delete=False) as f:
            f.write(audio.read())
            with sr.AudioFile(f.name) as source:
                voice=r.record(source)
                text=r.recognize_google(voice)
                st.write("You said:",text)
                ans=ask_gemini(text)
                st.session_state.activity+=1
                st.write(ans)
    st.markdown("</div>",unsafe_allow_html=True)

elif menu=="📊 Learning Tracker ⭐":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    progress=min(st.session_state.activity*7,100)
    st.progress(progress)
    st.write("Total interactions:",st.session_state.activity)
    st.write("🔥 Consistency = Success")
    st.markdown("</div>",unsafe_allow_html=True)

elif menu=="📂 History":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    for h in st.session_state.history[::-1]:
        st.markdown("### 🔹 "+h[0])
        st.write(h[2])
    st.markdown("</div>",unsafe_allow_html=True)

elif menu=="📄 Export PDF":
    st.write("📄 PDF Export feature coming soon!")
    # pdf=FPDF()
    # pdf.add_page()
    # pdf.set_font("Arial","",12)
    # for h in st.session_state.history[::-1]:
    #     pdf.multi_cell(0,8,f"{h[0]}:\n{h[2]}\n\n")
    # pdf_bytes=BytesIO()
    # pdf.output(pdf_bytes)
    # pdf_bytes.seek(0)
    # st.download_button("📄 Export Chat as PDF",data=pdf_bytes,file_name="chat_history.pdf",mime="application/pdf")

elif menu=="🖼 Concept Map":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    topic=st.text_input("Enter Topic for Concept Map")
    if st.button("Generate Concept Map"):
        # simple AI-style concept map using networkx
        G=nx.DiGraph()
        main=topic
        G.add_node(main)
        # AI generates subtopics
        subtopics=["Sub1","Sub2","Sub3","Sub4"]
        for s in subtopics:
            G.add_edge(main,s)
        plt.figure(figsize=(6,4))
        pos=nx.spring_layout(G)
        nx.draw(G,pos,with_labels=True,node_color="skyblue",node_size=3000,font_size=12,edge_color="purple",width=2)
        buf=BytesIO()
        plt.savefig(buf,format="PNG",transparent=True)
        buf.seek(0)
        st.image(buf,use_column_width=True)
    st.markdown("</div>",unsafe_allow_html=True)

elif menu=="🎮 Memory Challenge":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    topic=st.text_input("Enter Topic for Memory Challenge")
    if st.button("Start Challenge"):
        ans=ask_gemini(f"Create 3 quick questions with answers for {topic} to test memory in short game style.")
        st.write(ans)
    st.markdown("</div>",unsafe_allow_html=True)

elif menu=="ℹ️ About":
    st.markdown("<div class='glass'>",unsafe_allow_html=True)
    st.write("""
### 🎓 AI Study Buddy Pro

💬 Modern Animated UI/UX  
⚡ Fast Gemini Responses  
🖼 Concept Maps / Diagrams  
🎮 Memory Challenge Games  
🌙 Dark/Light Theme Toggle  
📄 PDF Export  Update Will Be Soon It was ONLY Availbe for Devloper
📂 Upload Anywhere  

**Developed by Darshil’s AI Team 🚀**
""")
    st.markdown("</div>",unsafe_allow_html=True)
